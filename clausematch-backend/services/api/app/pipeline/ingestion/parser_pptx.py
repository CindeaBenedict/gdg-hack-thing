from pathlib import Path
from pptx import Presentation


def parse_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text:
                texts.append(text)
        slides.append(f"Slide {i}:\n" + "\n".join(texts))
    return "\n\n".join(slides)


